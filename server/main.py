print("=== STARTING main.py ===", flush=True)
import uuid
import io
from typing import Optional, Union
from datetime import datetime

print("1. Importing FastAPI...", flush=True)
from fastapi import FastAPI, Depends, HTTPException, UploadFile, File
print("2. FastAPI imported. Importing CORS...", flush=True)
from fastapi.middleware.cors import CORSMiddleware
print("3. Importing SQLAlchemy...", flush=True)
from sqlalchemy.orm import Session
print("4. Importing db service...", flush=True)
from services.db import get_db, engine, Base
print("5. Importing rag_service...", flush=True)
from services.rag_service import retrieve_relevant_notes, build_rag_context, preload_model_async
print("6. Importing auth service...", flush=True)
from services.auth import (
    get_current_user, 
    get_current_user_optional,
    check_workspace_permission,
    get_user_workspaces
)
print("7. Importing models...", flush=True)
from models.note import Note
from models.workspace import Workspace
from models.user import User
from models.workspace_collaborator import (
    WorkspaceCollaborator, 
    PERMISSION_VIEWER, 
    PERMISSION_EDITOR, 
    PERMISSION_OWNER
)
from models import user, workspace, note
from models.workspace_collaborator import WorkspaceCollaborator as WC
print("8. Importing routers...", flush=True)
from routers.websocket_events import sio
from routers.auth import router as auth_router
from routers.collaborators import router as collaborators_router

import socketio
from pydantic import BaseModel

from openai import OpenAI
import os
from dotenv import load_dotenv
from re import finditer

load_dotenv()

print("Starting AI Colab Workspace API v2.0.0...", flush=True)
print("Initializing database...", flush=True)

try:
    Base.metadata.create_all(bind=engine)
    print("Database initialized successfully!", flush=True)
except Exception as e:
    # This is usually fine - tables already exist
    print(f"Note: Database tables already exist (this is OK): {type(e).__name__}", flush=True)

print("Setting up FastAPI application...", flush=True)

app = FastAPI(title="AI Colab Workspace API", version="2.0.0")

@app.on_event("startup")
async def startup_event():
    """Preload heavy models in background after server starts."""
    preload_model_async()

# CORS configuration - explicit origins required when using credentials
cors_origins = [
    "http://localhost:5173",
    "http://localhost:3000",
    "https://workspace.jacksmith.me",
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=cors_origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Include routers
app.include_router(auth_router)
app.include_router(collaborators_router)

openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

class ChatRequest(BaseModel):
    message: str
    workspace_id: str  # UUID as string
    conversation_history: list[dict] = []
    use_rag: bool = True

def extract_citations(message, relevant_notes):
    citation_pattern = r'\[Doc:\s*"([^"]+)"\]'
    matches = finditer(citation_pattern, message)

    title_to_note = {note["title"]: note["note_id"] for note in relevant_notes}

    citations = []
    for match in matches:
        cited_title = match.group(1)
        note_id = None
        for note_title, nid in title_to_note.items():
            if cited_title.lower() in note_title.lower() or note_title.lower() in cited_title.lower():
                note_id = nid
                break
        if note_id:
            citations.append({
                "note_id": note_id,
                "title": cited_title,
                "position": match.start(),
                "match_text": match.group(0)
            })

    return citations

@app.post("/ai/chat")
async def ai_chat(
    request: ChatRequest, 
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    # Convert workspace_id string to UUID
    ws_uuid = uuid.UUID(request.workspace_id)
    
    # Check permission (viewers can use AI)
    check_workspace_permission(db, current_user, ws_uuid, PERMISSION_VIEWER)
    
    try:
        relevant_notes = []
        rag_context = ""
        use_documents = False

        if request.use_rag:
            relevant_notes = retrieve_relevant_notes(
                query=request.message,
                workspace_id=ws_uuid,
                db=db,
                top_n=3
            )
            if relevant_notes and relevant_notes[0]["similarity"] > 0.15:
                rag_context = build_rag_context(relevant_notes)
                use_documents = True
        
        system_prompt = """
        You are an assistant that helps users work with their notes and documents in this workspace.

        You should:
        - Answer general questions using your own knowledge.
        - When a question is about the user’s documents, use the provided document context.
        - Keep answers concise, clear, and helpful.

        Formatting (Markdown):
        - Use **bold** for important terms or headings.
        - Use *italics* for subtle emphasis or titles.
        - Use `code` for inline code, variables, or technical identifiers.
        - Use bullet lists (- item) for unordered lists.
        - Use numbered lists (1. item) when order matters.
        - Use > for short callouts or quotes.
        """

        if use_documents and rag_context:
            system_prompt += f"""

            The following documents may be relevant to the user’s question:

            {rag_context}

            When you use information from these documents, follow these rules for citations:

            1. Citation format
            - Whenever you rely on a specific document, include a citation in this exact format:
                [Doc: "Document Title"]

            2. What requires a citation
            You must add a citation whenever you:
            - Quote text from a document.
            - Paraphrase or summarize its content.
            - List data taken from it (numbers, bullet points, etc.).
            - Analyse data that comes from it (e.g. means, ranges, trends).
            - Refer to the document by name (e.g. “Big Numbers”, “Who am I?”).
            - Describe what the document contains or shows.

            3. Where to place the citation
            - Prefer to introduce the document first, then describe its content.
            - Good pattern: [Doc: "Who am I?"] includes the statement "Hello, I am Shakespeare, that is a big name."
            - Good pattern: According to [Doc: "Meeting Notes"], the meeting is on Tuesday.
            - Good pattern: [Doc: "Big Numbers"] contains the numbers 2, 9, 1, 7, 5, ...
            - Avoid putting the citation as an afterthought at the very end of a sentence.

            4. Multiple citations
            - You can and should cite the same document more than once if you refer to it multiple times.
            - If you are describing data or analysis that clearly comes from one document, keep citing that document as needed so the source is obvious.

            5. General vs. document‑based answers
            - If the user’s question is general knowledge (not about their documents), answer normally and do not mention the documents.
            - If the question is clearly about the contents of the workspace (e.g. “Do I have any documents that mention Shakespeare?”, “What numbers appear in my files?”), base your answer on the documents above and use citations.

            Examples:

            - User: Do I have any documents that mention Shakespeare?
            Assistant: Yes, you do. [Doc: "Who am I?"] includes the statement "Hello, I am Shakespeare, that is a big name."

            - User: I have two files with numbers in them, find the files, list the numbers and then do statistical analysis.
            Assistant:
            - You have two documents with numbers:
                - [Doc: "Big Numbers"] contains 2, 9, 1, 7, 5, 10, 9, 7, 8, 9, 10, 9.
                - [Doc: "Really big numbers"] contains 18, 19, 36, 124, 9, 9, 4, 5, 4, 6.
            - [Doc: "Big Numbers"] has a mean of …
            - [Doc: "Really big numbers"] has a mean of …

            If you are unsure whether something comes from a document or not, either check the context carefully or say you are not certain, rather than guessing.
            """
        messages = [
            {"role": "system", "content": system_prompt}
        ]
        
        messages.extend(request.conversation_history)
        
        messages.append({"role": "user", "content": request.message})
        
        response = openai_client.chat.completions.create(
            model="gpt-4o-mini",
            messages=messages,
            max_tokens=1000,
            temperature=0.3
        )
        
        assistant_message = response.choices[0].message.content
        print(assistant_message)

        citations = extract_citations(assistant_message, relevant_notes)
        print(citations)
        
        source_data = [
            {
                "note_id": note["note_id"],
                "title": note["title"],
                "content": note["content"][:200],
                "similarity": note["similarity"],
                "created_at": note.get("created_at")
            }
            for note in relevant_notes
        ] if relevant_notes else []
        
        return {
            "message": assistant_message,
            "model": "gpt-4o-mini",
            "usage": {
                "prompt_tokens": response.usage.prompt_tokens,
                "completion_tokens": response.usage.completion_tokens,
                "total_tokens": response.usage.total_tokens
            },
            "sources": source_data,
            "citations": citations
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI chat error: {str(e)}")


# Supported file types for upload
SUPPORTED_TEXT_EXTENSIONS = {'.txt', '.md', '.markdown', '.json', '.csv', '.xml', '.html', '.htm', '.py', '.js', '.ts', '.jsx', '.tsx', '.css', '.yaml', '.yml', '.toml', '.ini', '.cfg', '.log', '.sql', '.sh', '.bat', '.ps1'}
SUPPORTED_DOCUMENT_EXTENSIONS = {'.pdf', '.docx', '.doc', '.rtf', '.odt', '.pptx', '.xlsx'}
ALL_SUPPORTED_EXTENSIONS = SUPPORTED_TEXT_EXTENSIONS | SUPPORTED_DOCUMENT_EXTENSIONS
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB for documents


def extract_text_from_pdf(content: bytes) -> str:
    """Extract text from PDF file."""
    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(content))
        text_parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
        return "\n\n".join(text_parts)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Could not extract text from PDF: {str(e)}")


def extract_text_from_docx(content: bytes) -> str:
    """Extract text from DOCX file."""
    try:
        import docx
        doc = docx.Document(io.BytesIO(content))
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_parts.append(cell.text)
        return "\n\n".join(text_parts)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Could not extract text from DOCX: {str(e)}")


def extract_text_from_pptx(content: bytes) -> str:
    """Extract text from PPTX file."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        return "\n\n".join(text_parts)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Could not extract text from PPTX: {str(e)}")


def extract_text_from_xlsx(content: bytes) -> str:
    """Extract text from XLSX file."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
        text_parts = []
        for sheet in wb.worksheets:
            text_parts.append(f"## Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    text_parts.append(row_text)
        return "\n".join(text_parts)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Could not extract text from XLSX: {str(e)}")


def extract_text_from_document(content: bytes, ext: str) -> str:
    """Extract text from document based on extension."""
    if ext == '.pdf':
        return extract_text_from_pdf(content)
    elif ext in {'.docx', '.doc'}:
        return extract_text_from_docx(content)
    elif ext == '.pptx':
        return extract_text_from_pptx(content)
    elif ext == '.xlsx':
        return extract_text_from_xlsx(content)
    elif ext in {'.rtf', '.odt'}:
        # For RTF and ODT, try basic text extraction or return error
        raise HTTPException(status_code=400, detail=f"File type {ext} support coming soon. Please convert to PDF or DOCX.")
    else:
        raise HTTPException(status_code=400, detail=f"Unsupported document type: {ext}")


@app.post("/workspaces/{workspace_id}/upload")
async def upload_file(
    workspace_id: str,
    file: UploadFile = File(...),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Upload a file and create a note from its content."""
    # Convert workspace_id to UUID
    try:
        ws_uuid = uuid.UUID(workspace_id)
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid workspace ID")
    
    # Check permission (need editor access to upload)
    check_workspace_permission(db, current_user, ws_uuid, PERMISSION_EDITOR)
    
    # Get file extension
    filename = file.filename or "uploaded_file"
    ext = os.path.splitext(filename)[1].lower()
    
    # Check file type
    if ext not in ALL_SUPPORTED_EXTENSIONS:
        raise HTTPException(
            status_code=400, 
            detail=f"Unsupported file type '{ext}'. Supported types: {', '.join(sorted(ALL_SUPPORTED_EXTENSIONS))}"
        )
    
    # Read file content
    try:
        content = await file.read()
        
        # Check file size
        if len(content) > MAX_FILE_SIZE:
            raise HTTPException(status_code=400, detail="File too large. Maximum size is 10MB.")
        
        # Extract text based on file type
        if ext in SUPPORTED_DOCUMENT_EXTENSIONS:
            # Document files - need special extraction
            text_content = extract_text_from_document(content, ext)
        else:
            # Text files - decode directly
            try:
                text_content = content.decode('utf-8')
            except UnicodeDecodeError:
                try:
                    text_content = content.decode('latin-1')
                except:
                    raise HTTPException(status_code=400, detail="Could not decode file. Please ensure it's a text file.")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading file: {str(e)}")
    
    # Create note from file
    title = os.path.splitext(filename)[0]  # Use filename without extension as title
    
    new_note = Note(
        workspace_id=ws_uuid,
        author_id=current_user.id,
        title=title,
        content=text_content
    )
    db.add(new_note)
    db.commit()
    db.refresh(new_note)
    
    # Emit socket event to notify other users
    from routers.websocket_events import sio
    await sio.emit('note_created', {
        'id': str(new_note.id),
        'workspace_id': str(new_note.workspace_id),
        'title': new_note.title,
        'content': new_note.content,
        'author_id': str(new_note.author_id),
        'created_at': new_note.created_at.isoformat() if new_note.created_at else None,
        'updated_at': new_note.updated_at.isoformat() if new_note.updated_at else None,
    }, room=f"workspace_{workspace_id}")
    
    return {
        "id": str(new_note.id),
        "title": new_note.title,
        "content": new_note.content,
        "filename": filename,
        "message": f"Successfully uploaded '{filename}' as a new note"
    }


@app.post("/workspaces/{workspace_id}/upload-multiple")
async def upload_multiple_files(
    workspace_id: str,
    files: list[UploadFile] = File(...),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Upload multiple files and create notes from their content."""
    # Convert workspace_id to UUID
    try:
        ws_uuid = uuid.UUID(workspace_id)
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid workspace ID")
    
    # Check permission (need editor access to upload)
    check_workspace_permission(db, current_user, ws_uuid, PERMISSION_EDITOR)
    
    results = []
    errors = []
    
    for file in files:
        filename = file.filename or "uploaded_file"
        ext = os.path.splitext(filename)[1].lower()
        
        # Check file type
        if ext not in ALL_SUPPORTED_EXTENSIONS:
            errors.append({"filename": filename, "error": f"Unsupported file type '{ext}'"})
            continue
        
        try:
            content = await file.read()
            
            # Check file size
            if len(content) > MAX_FILE_SIZE:
                errors.append({"filename": filename, "error": "File too large (max 10MB)"})
                continue
            
            # Extract text based on file type
            try:
                if ext in SUPPORTED_DOCUMENT_EXTENSIONS:
                    text_content = extract_text_from_document(content, ext)
                else:
                    try:
                        text_content = content.decode('utf-8')
                    except UnicodeDecodeError:
                        text_content = content.decode('latin-1')
            except HTTPException as he:
                errors.append({"filename": filename, "error": he.detail})
                continue
            except:
                errors.append({"filename": filename, "error": "Could not extract text from file"})
                continue
            
            # Create note
            title = os.path.splitext(filename)[0]
            new_note = Note(
                workspace_id=ws_uuid,
                author_id=current_user.id,
                title=title,
                content=text_content
            )
            db.add(new_note)
            db.commit()
            db.refresh(new_note)
            
            # Emit socket event
            from routers.websocket_events import sio
            await sio.emit('note_created', {
                'id': str(new_note.id),
                'workspace_id': str(new_note.workspace_id),
                'title': new_note.title,
                'content': new_note.content,
                'author_id': str(new_note.author_id),
                'created_at': new_note.created_at.isoformat() if new_note.created_at else None,
                'updated_at': new_note.updated_at.isoformat() if new_note.updated_at else None,
            }, room=f"workspace_{workspace_id}")
            
            results.append({
                "id": str(new_note.id),
                "title": new_note.title,
                "filename": filename
            })
        except Exception as e:
            errors.append({"filename": filename, "error": str(e)})
    
    return {
        "uploaded": results,
        "errors": errors,
        "message": f"Uploaded {len(results)} file(s)" + (f", {len(errors)} failed" if errors else "")
    }


@app.get("/ping")
def ping():
    return {"message": "pong", "version": "2.0.0"}


def serialize_workspace(item: Workspace, include_stats: bool = False, db: Session = None, user: User = None):
    data = {
        "id": item.id,
        "name": item.name,
        "description": item.description,
        "owner_id": item.owner_id,
        "is_shared": item.is_shared,
        "created_at": item.created_at.isoformat() if item.created_at else None,
        "updated_at": item.updated_at.isoformat() if item.updated_at else None,
    }
    
    if include_stats and db:
        # Count notes
        note_count = db.query(Note).filter(Note.workspace_id == item.id).count()
        data["note_count"] = note_count
        
        # Count members (owner + collaborators)
        member_count = db.query(WorkspaceCollaborator).filter(
            WorkspaceCollaborator.workspace_id == item.id,
            WorkspaceCollaborator.accepted_at.isnot(None)
        ).count() + 1  # +1 for owner
        data["member_count"] = member_count
        
        # Check user's role
        if user:
            if item.owner_id == user.id:
                data["role"] = "owner"
            else:
                collab = db.query(WorkspaceCollaborator).filter(
                    WorkspaceCollaborator.workspace_id == item.id,
                    WorkspaceCollaborator.user_id == user.id
                ).first()
                if collab:
                    if collab.permission_level == PERMISSION_EDITOR:
                        data["role"] = "editor"
                    else:
                        data["role"] = "viewer"
                else:
                    data["role"] = "none"
    
    return data


def serialize_note(item: Note):
    return {
        "id": item.id,
        "title": item.title,
        "content": item.content,
        "workspace_id": item.workspace_id,
        "author_id": item.author_id,
        "created_at": item.created_at.isoformat() if item.created_at else None,
        "updated_at": item.updated_at.isoformat() if item.updated_at else None,
    }


@app.get("/workspaces/")
async def list_workspaces(
    include_stats: bool = False,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """List all workspaces the current user has access to"""
    workspaces = get_user_workspaces(db, current_user)
    return [serialize_workspace(ws, include_stats=include_stats, db=db, user=current_user) for ws in workspaces]


class WorkspaceCreate(BaseModel):
    name: str
    description: Optional[str] = ""


class NoteCreate(BaseModel):
    title: str
    content: str = ""


class NoteUpdate(BaseModel):
    title: Optional[str] = None
    content: Optional[str] = None


@app.post("/workspaces/")
async def create_workspace(
    payload: WorkspaceCreate, 
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Create a new workspace"""
    workspace = Workspace(
        name=payload.name, 
        description=payload.description or "",
        owner_id=current_user.id
    )
    db.add(workspace)
    db.commit()
    db.refresh(workspace)
    return serialize_workspace(workspace)


@app.get("/workspaces/{workspace_id}")
async def get_workspace_detail(
    workspace_id: str, 
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Get workspace details"""
    ws_uuid = uuid.UUID(workspace_id)
    check_workspace_permission(db, current_user, ws_uuid, PERMISSION_VIEWER)
    
    workspace = db.query(Workspace).filter(Workspace.id == ws_uuid).first()
    if not workspace:
        raise HTTPException(status_code=404, detail="Workspace not found")
    
    return serialize_workspace(workspace, include_stats=True, db=db, user=current_user)


@app.delete("/workspaces/{workspace_id}")
async def delete_workspace(
    workspace_id: str,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Delete a workspace (owner only)"""
    ws_uuid = uuid.UUID(workspace_id)
    workspace = db.query(Workspace).filter(Workspace.id == ws_uuid).first()
    if not workspace:
        raise HTTPException(status_code=404, detail="Workspace not found")
    
    if workspace.owner_id != current_user.id:
        raise HTTPException(status_code=403, detail="Only the owner can delete this workspace")
    
    db.delete(workspace)
    db.commit()
    return {"message": "Workspace deleted"}


@app.get("/workspaces/{workspace_id}/notes/")
async def list_workspace_notes(
    workspace_id: str, 
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """List all notes in a workspace"""
    ws_uuid = uuid.UUID(workspace_id)
    check_workspace_permission(db, current_user, ws_uuid, PERMISSION_VIEWER)
    
    notes = (
        db.query(Note)
        .filter(Note.workspace_id == ws_uuid)
        .order_by(Note.updated_at.desc())
        .all()
    )
    return [serialize_note(note) for note in notes]


@app.post("/workspaces/{workspace_id}/notes/")
async def create_note(
    workspace_id: str,
    payload: NoteCreate,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Create a new note in a workspace"""
    ws_uuid = uuid.UUID(workspace_id)
    check_workspace_permission(db, current_user, ws_uuid, PERMISSION_EDITOR)
    
    note = Note(
        title=payload.title,
        content=payload.content,
        workspace_id=ws_uuid,
        author_id=current_user.id
    )
    db.add(note)
    db.commit()
    db.refresh(note)
    return serialize_note(note)


@app.get("/notes/{note_id}")
async def get_note(
    note_id: int,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Get a specific note"""
    note = db.query(Note).filter(Note.id == note_id).first()
    if not note:
        raise HTTPException(status_code=404, detail="Note not found")
    
    check_workspace_permission(db, current_user, note.workspace_id, PERMISSION_VIEWER)
    
    return serialize_note(note)


@app.put("/notes/{note_id}")
async def update_note(
    note_id: int,
    payload: NoteUpdate,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Update a note"""
    note = db.query(Note).filter(Note.id == note_id).first()
    if not note:
        raise HTTPException(status_code=404, detail="Note not found")
    
    check_workspace_permission(db, current_user, note.workspace_id, PERMISSION_EDITOR)
    
    if payload.title is not None:
        note.title = payload.title
    if payload.content is not None:
        note.content = payload.content
    
    db.commit()
    db.refresh(note)
    return serialize_note(note)


@app.delete("/notes/{note_id}")
async def delete_note(
    note_id: int,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Delete a note (owner only)"""
    note = db.query(Note).filter(Note.id == note_id).first()
    if not note:
        raise HTTPException(status_code=404, detail="Note not found")
    
    workspace = db.query(Workspace).filter(Workspace.id == note.workspace_id).first()
    if workspace.owner_id != current_user.id:
        raise HTTPException(status_code=403, detail="Only workspace owner can delete notes")
    
    db.delete(note)
    db.commit()
    return {"message": "Note deleted"}


@app.get("/dashboard/stats")
async def get_dashboard_stats(
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Get dashboard statistics for current user"""
    workspaces = get_user_workspaces(db, current_user)
    workspace_count = len(workspaces)
    
    workspace_ids = [w.id for w in workspaces]
    note_count = db.query(Note).filter(Note.workspace_id.in_(workspace_ids)).count() if workspace_ids else 0
    
    # Wrap collaborator queries in try-except since table might not exist
    pending_invitations = 0
    collaborator_count = 0
    try:
        pending_invitations = db.query(WorkspaceCollaborator).filter(
            WorkspaceCollaborator.user_id == current_user.id,
            WorkspaceCollaborator.accepted_at.is_(None)
        ).count()
        
        for ws in workspaces:
            if ws.owner_id == current_user.id:
                collabs = db.query(WorkspaceCollaborator).filter(
                    WorkspaceCollaborator.workspace_id == ws.id,
                    WorkspaceCollaborator.accepted_at.isnot(None)
                ).count()
                collaborator_count += collabs
    except Exception:
        # Table might not exist yet
        db.rollback()
        pending_invitations = 0
        collaborator_count = 0
    
    return {
        "workspace_count": workspace_count,
        "note_count": note_count,
        "collaborator_count": collaborator_count,
        "pending_invitations": pending_invitations
    }

socket_app = socketio.ASGIApp(sio, other_asgi_app=app)

if __name__ == "__main__":
    import uvicorn
    print("Starting uvicorn server on http://0.0.0.0:8000 ...")
    uvicorn.run("main:socket_app", host="0.0.0.0", port=8000, reload=True)